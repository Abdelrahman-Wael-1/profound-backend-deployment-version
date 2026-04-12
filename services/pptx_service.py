"""
pptx_service.py
Professional university lecture PowerPoint generator.

Layout per content slide:
  ┌─────────────────────────────────────────────────────┐
  │ ▌  SLIDE TITLE                              [  N  ] │
  │ ═══════════════════════════════════════════════════ │
  │                                                     │
  │  ● Headline One                                     │
  │    Detail explanation for headline one in 2-3       │
  │    sentences giving depth and context.              │
  │                                                     │
  │  ● Headline Two                                     │
  │    Detail explanation for headline two…             │
  │                                                     │
  │  ● Headline Three                                   │
  │    Detail explanation for headline three…           │
  │                                                     │
  │ ┌─────────────────────────────────────────────────┐ │
  │ │ EXAMPLE  Concrete real-world worked example…    │ │
  │ └─────────────────────────────────────────────────┘ │
  └─────────────────────────────────────────────────────┘
"""

import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Theme colour palettes (no image files — pure colour) ──────────
THEMES = {
    "Modern Minimalist": {
        "bg":         (0xFF, 0xFF, 0xFF),
        "accent":     (0x4F, 0x46, 0xE5),   # indigo-600
        "accent2":    (0x7C, 0x3A, 0xED),   # violet-600
        "title_txt":  (0x1E, 0x1B, 0x4B),   # indigo-950
        "headline":   (0x1E, 0x1B, 0x4B),   # same deep indigo
        "detail_txt": (0x37, 0x41, 0x51),   # gray-700
        "ex_bg":      (0xEC, 0xFD, 0xF5),   # emerald-50
        "ex_label":   (0x06, 0x5F, 0x46),   # emerald-900
        "ex_txt":     (0x06, 0x5F, 0x46),
        "dot":        (0x4F, 0x46, 0xE5),
        "bar":        (0x4F, 0x46, 0xE5),
        "divider":    (0xC7, 0xD2, 0xFE),   # indigo-200
        "badge_bg":   (0x4F, 0x46, 0xE5),
        "badge_txt":  (0xFF, 0xFF, 0xFF),
    },
    "Dark Mode Tech": {
        "bg":         (0x0F, 0x17, 0x2A),
        "accent":     (0x38, 0xBD, 0xF8),   # sky-400
        "accent2":    (0x06, 0xB6, 0xD4),
        "title_txt":  (0xF8, 0xFA, 0xFC),
        "headline":   (0xE2, 0xE8, 0xF0),
        "detail_txt": (0x94, 0xA3, 0xB8),   # slate-400
        "ex_bg":      (0x0C, 0x2A, 0x2A),
        "ex_label":   (0x34, 0xD3, 0x99),
        "ex_txt":     (0x34, 0xD3, 0x99),
        "dot":        (0x38, 0xBD, 0xF8),
        "bar":        (0x38, 0xBD, 0xF8),
        "divider":    (0x1E, 0x40, 0x4F),
        "badge_bg":   (0x38, 0xBD, 0xF8),
        "badge_txt":  (0x0F, 0x17, 0x2A),
    },
    "Classic Academic": {
        "bg":         (0xFD, 0xFB, 0xF7),
        "accent":     (0x80, 0x00, 0x00),   # maroon
        "accent2":    (0xB8, 0x5C, 0x38),
        "title_txt":  (0x3B, 0x0A, 0x0A),
        "headline":   (0x3B, 0x0A, 0x0A),
        "detail_txt": (0x1E, 0x1E, 0x1E),
        "ex_bg":      (0xF0, 0xF7, 0xEE),
        "ex_label":   (0x1A, 0x47, 0x2A),
        "ex_txt":     (0x1A, 0x47, 0x2A),
        "dot":        (0x80, 0x00, 0x00),
        "bar":        (0x80, 0x00, 0x00),
        "divider":    (0xD9, 0xC5, 0xB2),
        "badge_bg":   (0x80, 0x00, 0x00),
        "badge_txt":  (0xFF, 0xFF, 0xFF),
    },
    "Vibrant Creative": {
        "bg":         (0xFF, 0xFF, 0xFF),
        "accent":     (0xF9, 0x73, 0x16),   # orange-500
        "accent2":    (0xEC, 0x48, 0x99),
        "title_txt":  (0x43, 0x14, 0x07),
        "headline":   (0x43, 0x14, 0x07),
        "detail_txt": (0x1C, 0x1C, 0x1E),
        "ex_bg":      (0xFD, 0xF2, 0xF8),
        "ex_label":   (0x70, 0x1A, 0x75),
        "ex_txt":     (0x70, 0x1A, 0x75),
        "dot":        (0xF9, 0x73, 0x16),
        "bar":        (0xF9, 0x73, 0x16),
        "divider":    (0xFE, 0xD7, 0xAA),
        "badge_bg":   (0xF9, 0x73, 0x16),
        "badge_txt":  (0xFF, 0xFF, 0xFF),
    },
}

_W = Inches(13.333)
_H = Inches(7.5)


def _rgb(t):
    return RGBColor(*t)


def _solid_rect(slide, left, top, width, height, color):
    """Add a solid-fill rectangle with no border."""
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(color)
    s.line.fill.background()
    return s


def _textbox(slide, left, top, width, height, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    return tf


def _para(tf, text, size, bold, color, align=PP_ALIGN.LEFT, space_before=0, space_after=0, italic=False):
    """Append a paragraph to an existing text frame. First empty para is reused."""
    if tf.paragraphs and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text          = str(text)
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = _rgb(color)
    p.alignment       = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after  = Pt(space_after)
    return p


def create_pptx(data: dict) -> io.BytesIO:
    theme_name  = data.get("theme", "Modern Minimalist")
    c           = THEMES.get(theme_name, THEMES["Modern Minimalist"])
    course_lvl  = data.get("course_level", "")

    prs             = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H

    slides_data = data.get("slides", [])
    if not slides_data:
        raise ValueError("No slides data provided")

    for idx, sd in enumerate(slides_data):
        slide      = prs.slides.add_slide(prs.slide_layouts[6])   # blank
        title      = str(sd.get("title", "")).strip()
        is_title   = (idx == 0)
        is_ref     = any(k in title.lower() for k in ("reference", "further", "resource"))

        # ── Slide background ──────────────────────────────────────
        _solid_rect(slide, 0, 0, _W, _H, c["bg"])

        # ── Left accent bar ───────────────────────────────────────
        _solid_rect(slide, 0, 0, Inches(0.08), _H, c["bar"])

        # ══════════════════════════════════════════════════════════
        # TITLE SLIDE
        # ══════════════════════════════════════════════════════════
        if is_title:
            # Large decorative circle top-right
            cir = slide.shapes.add_shape(9, Inches(8.8), Inches(-1.8), Inches(5.5), Inches(5.5))
            cir.fill.solid(); cir.fill.fore_color.rgb = _rgb(c["accent"]); cir.line.fill.background()
            cir2 = slide.shapes.add_shape(9, Inches(10.0), Inches(4.2), Inches(4.0), Inches(4.0))
            cir2.fill.solid(); cir2.fill.fore_color.rgb = _rgb(c["accent2"]); cir2.line.fill.background()

            # Title
            tf = _textbox(slide, Inches(1.1), Inches(1.8), Inches(9.0), Inches(2.8))
            _para(tf, title,    48, True,  c["title_txt"], PP_ALIGN.LEFT)
            if course_lvl:
                _para(tf, course_lvl, 20, False, c["accent"],    PP_ALIGN.LEFT, space_before=8)

            # Underline accent
            _solid_rect(slide, Inches(1.1), Inches(4.75), Inches(3.0), Inches(0.07), c["accent"])
            continue

        # ── Slide number badge ────────────────────────────────────
        badge_w = Inches(0.55)
        badge_h = Inches(0.32)
        _solid_rect(slide, _W - badge_w - Inches(0.12), _H - badge_h - Inches(0.1), badge_w, badge_h, c["badge_bg"])
        btf = _textbox(slide, _W - badge_w - Inches(0.10), _H - badge_h - Inches(0.08), badge_w, badge_h)
        _para(btf, str(idx + 1), 10, True, c["badge_txt"], PP_ALIGN.CENTER)

        # ── Title ─────────────────────────────────────────────────
        ttf = _textbox(slide, Inches(0.22), Inches(0.10), Inches(12.8), Inches(0.80))
        _para(ttf, title, 26, True, c["title_txt"])

        # ── Title underline ───────────────────────────────────────
        _solid_rect(slide, Inches(0.22), Inches(0.93), Inches(12.8), Inches(0.035), c["accent"])

        # ══════════════════════════════════════════════════════════
        # REFERENCE / RESOURCE SLIDE
        # ══════════════════════════════════════════════════════════
        if is_ref:
            points = sd.get("points", [])
            content_tf = _textbox(slide, Inches(0.35), Inches(1.05), Inches(12.6), Inches(5.8))
            for pt in points:
                headline = str(pt.get("headline", "")).strip()
                detail   = str(pt.get("detail",   "")).strip()
                # Detect URL in detail field
                is_url = detail.startswith("http")
                _para(content_tf, f"  {headline}", 16, True,  c["headline"],   space_before=14)
                if detail:
                    _para(content_tf, f"  {detail}", 13, False, c["accent"] if is_url else c["detail_txt"], space_before=2)

            # Example as note at bottom
            example = str(sd.get("example", "")).strip()
            if example:
                _solid_rect(slide, Inches(0.22), Inches(6.6), Inches(12.8), Inches(0.72), c["ex_bg"])
                etf = _textbox(slide, Inches(0.40), Inches(6.65), Inches(12.5), Inches(0.62))
                _para(etf, example, 11, False, c["ex_txt"])
            continue

        # ══════════════════════════════════════════════════════════
        # CONTENT SLIDE — main layout
        # ══════════════════════════════════════════════════════════
        points  = sd.get("points", [])
        example = str(sd.get("example", "")).strip()
        img_sug = sd.get("image_suggestion")
        notes   = str(sd.get("speaker_notes", "")).strip()

        # Decide vertical space allocation
        has_example = bool(example)
        ex_h        = Inches(1.10)
        ex_top      = _H - ex_h - Inches(0.12)
        content_top = Inches(1.05)
        content_h   = (ex_top - content_top - Inches(0.12)) if has_example else (_H - content_top - Inches(0.15))

        # ── Points section ────────────────────────────────────────
        # We render each point as: bullet dot | HEADLINE (bold) then detail below indented
        point_tf = _textbox(slide, Inches(0.22), content_top, Inches(12.8), content_h)

        for i, pt in enumerate(points):
            headline = str(pt.get("headline", "")).strip()
            detail   = str(pt.get("detail",   "")).strip()

            space_top = 10 if i == 0 else 18   # more breathing room between points

            # ── Headline row: accent dot + bold headline ──────────
            # We render "●  HEADLINE" as one paragraph with two runs
            if point_tf.paragraphs and point_tf.paragraphs[0].text == "":
                h_para = point_tf.paragraphs[0]
            else:
                h_para = point_tf.add_paragraph()

            # Dot run (accent colour)
            dot_run            = h_para.add_run()
            dot_run.text       = "●  "
            dot_run.font.size  = Pt(13)
            dot_run.font.bold  = True
            dot_run.font.color.rgb = _rgb(c["dot"])

            # Headline run
            hl_run             = h_para.add_run()
            hl_run.text        = headline
            hl_run.font.size   = Pt(18)
            hl_run.font.bold   = True
            hl_run.font.color.rgb = _rgb(c["headline"])
            h_para.space_before = Pt(space_top)
            h_para.space_after  = Pt(2)

            # ── Detail paragraph (indented, smaller, normal weight) ─
            if detail:
                d_para = point_tf.add_paragraph()
                d_run            = d_para.add_run()
                d_run.text       = "     " + detail   # indent with spaces
                d_run.font.size  = Pt(13)
                d_run.font.bold  = False
                d_run.font.color.rgb = _rgb(c["detail_txt"])
                d_para.space_before = Pt(1)
                d_para.space_after  = Pt(2)
                d_para.alignment    = PP_ALIGN.LEFT

        # ── Thin separator before example ─────────────────────────
        if has_example:
            _solid_rect(slide, Inches(0.22), ex_top - Inches(0.08), Inches(12.8), Inches(0.025), c["divider"])

        # ── Example box ───────────────────────────────────────────
        if has_example:
            _solid_rect(slide, Inches(0.22), ex_top, Inches(12.8), ex_h, c["ex_bg"])

            ex_tf = _textbox(slide, Inches(0.40), ex_top + Inches(0.07), Inches(12.5), ex_h - Inches(0.1))

            # "EXAMPLE" label run
            label_para       = ex_tf.paragraphs[0]
            label_run        = label_para.add_run()
            label_run.text   = "EXAMPLE   "
            label_run.font.size  = Pt(10)
            label_run.font.bold  = True
            label_run.font.color.rgb = _rgb(c["ex_label"])

            # Example text run (same paragraph — label + text on one line for compactness)
            text_run         = label_para.add_run()
            text_run.text    = example
            text_run.font.size = Pt(12)
            text_run.font.bold = False
            text_run.font.color.rgb = _rgb(c["ex_txt"])

        # ── Image suggestion (subtle footer tag) ──────────────────
        if img_sug and str(img_sug).lower() not in ("null", "none", ""):
            img_tf = _textbox(slide, Inches(0.22), _H - Inches(0.28), Inches(8.0), Inches(0.25))
            _para(img_tf, f"🖼  {img_sug}", 8, False, c["accent"], italic=True)

        # ── Speaker notes ─────────────────────────────────────────
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    stream = io.BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream